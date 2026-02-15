from __future__ import annotations

import math
import os
import re
from collections import Counter
from datetime import datetime
from io import BytesIO

import pandas as pd
from docx import Document
from openpyxl import Workbook
from openpyxl.chart import BarChart, PieChart, Reference
from openpyxl.styles import Font, PatternFill
from pypdf import PdfReader
from pptx import Presentation
from .ai_runtime import semantic_key_points

EXCEL_MAX_ROWS = 1_048_576
MAX_DATA_ROWS_PER_SHEET = EXCEL_MAX_ROWS - 1
AUTO_WIDTH_SCAN_LIMIT = 250
DEFAULT_CLEANED_EXPORT_MAX_ROWS = 50_000
DEFAULT_ANALYSIS_SAMPLE_MAX_ROWS = 200_000
TYPE_INFERENCE_SAMPLE_ROWS = 2_000
DEFAULT_MAX_PROCESS_ROWS = 300_000


def _excel_value(value):
    if pd.isna(value):
        return ""
    if isinstance(value, pd.Timestamp):
        value = value.to_pydatetime()
    if isinstance(value, datetime):
        return value.replace(tzinfo=None)
    return value


def _apply_header_style(worksheet):
    for cell in worksheet[1]:
        cell.font = Font(bold=True)
        cell.fill = PatternFill(start_color="DDEBFF", end_color="DDEBFF", fill_type="solid")


def _autofit_columns(worksheet, max_scan_rows=AUTO_WIDTH_SCAN_LIMIT):
    max_row = min(worksheet.max_row, max_scan_rows)
    for col_idx, column_cells in enumerate(worksheet.iter_cols(min_row=1, max_row=max_row), start=1):
        max_len = max(len(str(cell.value)) if cell.value is not None else 0 for cell in column_cells)
        column_letter = worksheet.cell(row=1, column=col_idx).column_letter
        worksheet.column_dimensions[column_letter].width = min(max(max_len + 2, 10), 54)


def _write_table(worksheet, headers, rows):
    worksheet.append(headers)
    _apply_header_style(worksheet)
    written = 0
    for row in rows:
        if worksheet.max_row >= EXCEL_MAX_ROWS:
            break
        worksheet.append(row)
        written += 1
    _autofit_columns(worksheet)
    return written


def _write_dataframe_paginated(workbook: Workbook, base_sheet_name: str, dataframe: pd.DataFrame) -> int:
    if dataframe.empty:
        worksheet = workbook.create_sheet(base_sheet_name)
        _write_table(worksheet, list(dataframe.columns), [])
        return 1

    total_rows = len(dataframe)
    sheet_count = max(1, math.ceil(total_rows / MAX_DATA_ROWS_PER_SHEET))
    for index in range(sheet_count):
        start = index * MAX_DATA_ROWS_PER_SHEET
        end = min((index + 1) * MAX_DATA_ROWS_PER_SHEET, total_rows)
        title = f"{base_sheet_name}_{index + 1}" if sheet_count > 1 else base_sheet_name
        worksheet = workbook.create_sheet(title[:31])
        worksheet.append(list(dataframe.columns))
        _apply_header_style(worksheet)
        chunk = dataframe.iloc[start:end]
        for row in chunk.itertuples(index=False, name=None):
            worksheet.append([_excel_value(value) for value in row])
        _autofit_columns(worksheet)
    return sheet_count


def _flatten_columns(columns):
    flattened = []
    for col in columns:
        if isinstance(col, tuple):
            parts = [str(part) for part in col if str(part)]
            flattened.append(" | ".join(parts))
        else:
            flattened.append(str(col))
    return flattened


def _infer_series_type(series: pd.Series, sample_rows: int = TYPE_INFERENCE_SAMPLE_ROWS) -> str:
    sample = series.dropna()
    if sample.empty:
        return "object"
    if len(sample) > sample_rows:
        sample = sample.sample(n=sample_rows, random_state=42)

    numeric = pd.to_numeric(sample, errors="coerce")
    numeric_ratio = float(numeric.notna().mean()) if len(sample) else 0.0
    if numeric_ratio >= 0.95:
        return "numeric"

    parsed_dates = pd.to_datetime(sample, errors="coerce")
    date_ratio = float(parsed_dates.notna().mean()) if len(sample) else 0.0
    if date_ratio >= 0.95:
        return "datetime"

    return "object"


def _load_business_dataframe(uploaded_file, filename: str) -> pd.DataFrame:
    lower_name = filename.lower()
    if lower_name.endswith(".csv"):
        return pd.read_csv(uploaded_file)
    if lower_name.endswith(".xlsx") or lower_name.endswith(".xls"):
        return pd.read_excel(uploaded_file)
    raise ValueError("Only .xlsx, .xls, and .csv files are supported for data analysis.")


def analyze_business_data(uploaded_file, filename: str) -> tuple[dict, bytes]:
    raw_df = _load_business_dataframe(uploaded_file, filename)
    if raw_df.empty:
        raise ValueError("Uploaded dataset is empty.")

    original_shape = raw_df.shape
    raw_df = raw_df.dropna(how="all").copy()
    raw_df.columns = [str(col).strip() if str(col).strip() else f"column_{idx}" for idx, col in enumerate(raw_df.columns)]
    raw_df = raw_df.loc[:, ~raw_df.columns.duplicated()].copy()

    try:
        max_process_rows = int(os.getenv("OFFICE_MAX_PROCESS_ROWS", str(DEFAULT_MAX_PROCESS_ROWS)))
    except ValueError:
        max_process_rows = DEFAULT_MAX_PROCESS_ROWS
    max_process_rows = max(50_000, max_process_rows)

    large_dataset_mode = len(raw_df) > max_process_rows
    if large_dataset_mode:
        df = raw_df.sample(n=max_process_rows, random_state=42).copy()
    else:
        df = raw_df.copy()
    processing_input_rows = int(len(df))

    for col in df.columns:
        if df[col].dtype != "object":
            continue
        inferred = _infer_series_type(df[col])
        if inferred == "numeric":
            df[col] = pd.to_numeric(df[col], errors="coerce")
        elif inferred == "datetime":
            df[col] = pd.to_datetime(df[col], errors="coerce")

    numeric_cols = df.select_dtypes(include=["number"]).columns.tolist()
    datetime_cols = df.select_dtypes(include=["datetime64[ns]", "datetime64[ns, UTC]"]).columns.tolist()
    categorical_cols = [c for c in df.columns if c not in numeric_cols and c not in datetime_cols]

    missing_by_column_before = df.isna().sum().sort_values(ascending=False)
    missing_before = int(missing_by_column_before.sum())
    for col in numeric_cols:
        if df[col].isna().any():
            df[col] = df[col].fillna(df[col].median())
    for col in categorical_cols:
        if df[col].isna().any():
            mode = df[col].mode()
            fill_value = mode.iloc[0] if not mode.empty else "unknown"
            df[col] = df[col].fillna(fill_value)
    for col in datetime_cols:
        if df[col].isna().any():
            df[col] = df[col].fillna(method="ffill").fillna(method="bfill")

    duplicate_rows = int(df.duplicated().sum())
    df = df.drop_duplicates()
    rows_removed = int(processing_input_rows - len(df))

    try:
        analysis_sample_limit = int(
            os.getenv("OFFICE_ANALYSIS_SAMPLE_MAX_ROWS", str(DEFAULT_ANALYSIS_SAMPLE_MAX_ROWS))
        )
    except ValueError:
        analysis_sample_limit = DEFAULT_ANALYSIS_SAMPLE_MAX_ROWS
    analysis_sample_limit = max(50_000, analysis_sample_limit)
    analysis_df = df if len(df) <= analysis_sample_limit else df.sample(n=analysis_sample_limit, random_state=42)

    outlier_total = 0
    outlier_details = []
    for col in numeric_cols:
        series = analysis_df[col].dropna()
        if len(series) < 4:
            continue
        q1 = series.quantile(0.25)
        q3 = series.quantile(0.75)
        iqr = q3 - q1
        if iqr == 0:
            continue
        lower = q1 - 1.5 * iqr
        upper = q3 + 1.5 * iqr
        outliers = int(((analysis_df[col] < lower) | (analysis_df[col] > upper)).sum())
        outlier_total += outliers
        outlier_details.append([col, outliers, float(lower), float(upper)])

    pivot1 = None
    pivot2 = None
    if categorical_cols and numeric_cols:
        pivot1 = pd.pivot_table(
            analysis_df,
            index=categorical_cols[0],
            values=numeric_cols[0],
            aggfunc=["sum", "mean", "count"],
            fill_value=0,
        )
    if len(categorical_cols) > 1:
        pivot2 = pd.pivot_table(
            analysis_df,
            index=categorical_cols[0],
            columns=categorical_cols[1],
            values=numeric_cols[0] if numeric_cols else categorical_cols[0],
            aggfunc="count",
            fill_value=0,
        )

    summary = {
        "filename": filename,
        "rows_uploaded": int(original_shape[0]),
        "columns_uploaded": int(original_shape[1]),
        "rows_profiled": processing_input_rows,
        "rows_skipped_for_profiling": int(max(original_shape[0] - processing_input_rows, 0)),
        "rows_after_cleaning": int(len(df)),
        "rows_removed": rows_removed,
        "missing_cells_filled": missing_before,
        "duplicate_rows_removed": duplicate_rows,
        "numeric_columns": numeric_cols,
        "categorical_columns": categorical_cols,
        "datetime_columns": datetime_cols,
        "outlier_count": outlier_total,
        "analysis_sample_rows": int(len(analysis_df)),
        "large_dataset_mode": large_dataset_mode,
        "generated_at": datetime.utcnow().isoformat(),
    }
    try:
        cleaned_export_limit = int(os.getenv("OFFICE_CLEANED_EXPORT_MAX_ROWS", str(DEFAULT_CLEANED_EXPORT_MAX_ROWS)))
    except ValueError:
        cleaned_export_limit = DEFAULT_CLEANED_EXPORT_MAX_ROWS
    cleaned_export_limit = max(10_000, cleaned_export_limit)
    export_df = df if len(df) <= cleaned_export_limit else df.head(cleaned_export_limit)
    summary["cleaned_rows_exported"] = int(len(export_df))
    summary["cleaned_rows_truncated"] = int(max(len(df) - len(export_df), 0))

    workbook = Workbook()
    workbook.remove(workbook.active)

    ws_dashboard = workbook.create_sheet("Dashboard")
    _write_table(
        ws_dashboard,
        ["Metric", "Value"],
        [
            ["Rows Uploaded", summary["rows_uploaded"]],
            ["Columns Uploaded", summary["columns_uploaded"]],
            ["Rows Profiled", summary["rows_profiled"]],
            ["Rows After Cleaning", summary["rows_after_cleaning"]],
            ["Rows Removed", summary["rows_removed"]],
            ["Missing Cells Filled", summary["missing_cells_filled"]],
            ["Duplicate Rows Removed", summary["duplicate_rows_removed"]],
            ["Outlier Count", summary["outlier_count"]],
        ],
    )
    ws_dashboard["D1"] = "Analyst Workflow"
    ws_dashboard["D2"] = "Load -> Profile -> Clean -> Outlier Scan -> Pivot -> Visualize"
    ws_dashboard["D1"].font = Font(bold=True)

    ws_profile = workbook.create_sheet("Column_Profile")
    profile_rows = []
    for col in df.columns:
        profile_rows.append(
            [
                col,
                str(df[col].dtype),
                int(df[col].isna().sum()),
                int(df[col].nunique(dropna=True)),
                str(df[col].head(1).iloc[0]) if len(df[col]) else "",
            ]
        )
    _write_table(ws_profile, ["Column", "DType", "Missing", "Distinct", "Sample"], profile_rows)

    ws_missing = workbook.create_sheet("Missing_Before_Clean")
    _write_table(
        ws_missing,
        ["Column", "Missing Cells"],
        [[col, int(val)] for col, val in missing_by_column_before.items()],
    )

    cleaned_sheet_count = _write_dataframe_paginated(workbook, "Cleaned_Data", export_df)
    summary["cleaned_data_sheets"] = cleaned_sheet_count

    if outlier_details:
        ws_outliers = workbook.create_sheet("Outliers")
        _write_table(ws_outliers, ["Column", "Outlier Count", "Lower Bound", "Upper Bound"], outlier_details)

    if pivot1 is not None:
        ws_pivot1 = workbook.create_sheet("Pivot_1")
        p1 = pivot1.reset_index()
        _write_table(ws_pivot1, _flatten_columns(list(p1.columns)), p1.values.tolist())

        chart = BarChart()
        chart.title = f"{p1.columns[0]} vs {p1.columns[1]}"
        data_ref = Reference(ws_pivot1, min_col=2, min_row=1, max_col=2, max_row=ws_pivot1.max_row)
        cats_ref = Reference(ws_pivot1, min_col=1, min_row=2, max_row=ws_pivot1.max_row)
        chart.add_data(data_ref, titles_from_data=True)
        chart.set_categories(cats_ref)
        chart.height = 7
        chart.width = 11
        ws_dashboard.add_chart(chart, "A10")

    if pivot2 is not None:
        ws_pivot2 = workbook.create_sheet("Pivot_2")
        p2 = pivot2.reset_index()
        _write_table(ws_pivot2, _flatten_columns(list(p2.columns)), p2.values.tolist())

    if numeric_cols:
        ws_numeric = workbook.create_sheet("Numeric_Stats")
        numeric_stats = analysis_df[numeric_cols].describe().transpose().reset_index()
        _write_table(
            ws_numeric,
            _flatten_columns(list(numeric_stats.columns)),
            numeric_stats.fillna("").values.tolist(),
        )

        if len(numeric_cols) > 1:
            corr = analysis_df[numeric_cols].corr(numeric_only=True).reset_index()
            ws_corr = workbook.create_sheet("Correlation")
            _write_table(ws_corr, _flatten_columns(list(corr.columns)), corr.fillna("").values.tolist())

    if categorical_cols:
        ws_cat = workbook.create_sheet("Top_Categories")
        category_rows = []
        for column in categorical_cols[:5]:
            counts = analysis_df[column].astype(str).value_counts().head(15)
            total = max(int(counts.sum()), 1)
            for idx, val in counts.items():
                category_rows.append([column, idx, int(val), round((val / total) * 100, 2)])
        _write_table(ws_cat, ["Column", "Category", "Count", "Share %"], category_rows)

        top_counts = analysis_df[categorical_cols[0]].value_counts().head(10)
        pie = PieChart()
        pie.title = f"Top {categorical_cols[0]}"
        top_for_chart = [["Category", "Count"]]
        top_for_chart.extend([[idx, int(val)] for idx, val in top_counts.items()])
        ws_top = workbook.create_sheet("Top_Category_Chart")
        _write_table(ws_top, top_for_chart[0], top_for_chart[1:])
        data = Reference(ws_top, min_col=2, min_row=1, max_row=ws_top.max_row)
        labels = Reference(ws_top, min_col=1, min_row=2, max_row=ws_top.max_row)
        pie.add_data(data, titles_from_data=True)
        pie.set_categories(labels)
        pie.height = 7
        pie.width = 9
        ws_dashboard.add_chart(pie, "M10")

    ws_notes = workbook.create_sheet("Analyst_Notes")
    _write_table(
        ws_notes,
        ["Note"],
        [
            ["Open Pivot sheets in desktop Excel to add slicers and timeline controls."],
            ["The workbook is generated from uploaded business data, not application task records."],
            [f"Large cleaned datasets are split across {cleaned_sheet_count} sheet(s) to respect Excel row limits."],
            [f"Cleaned row export capped at {summary['cleaned_rows_exported']} rows for performance."],
            [f"Advanced stats/pivots computed on a representative sample of {summary['analysis_sample_rows']} rows for speed."],
            [f"Large dataset mode: {'Enabled' if large_dataset_mode else 'Disabled'}."],
        ],
    )

    output = BytesIO()
    workbook.save(output)
    output.seek(0)
    return summary, output.getvalue()


def _extract_keywords(text: str, limit: int = 8) -> list[str]:
    tokens = re.findall(r"[A-Za-z]{4,}", text.lower())
    stop_words = {"this", "that", "with", "from", "have", "will", "would", "about", "there", "were", "been"}
    counts = Counter(word for word in tokens if word not in stop_words)
    return [word for word, _ in counts.most_common(limit)]


def extract_document_text(uploaded_file, filename: str) -> str:
    name = filename.lower()
    if name.endswith(".txt") or name.endswith(".md"):
        return uploaded_file.read().decode("utf-8", errors="ignore")
    if name.endswith(".docx"):
        doc = Document(uploaded_file)
        return "\n".join(p.text.strip() for p in doc.paragraphs if p.text.strip())
    if name.endswith(".pdf"):
        reader = PdfReader(uploaded_file)
        chunks = []
        for page in reader.pages:
            chunks.append(page.extract_text() or "")
        return "\n".join(chunks).strip()
    raise ValueError("Supported document formats: .txt, .md, .docx, .pdf")


def build_powerpoint_report(source_name: str, text: str) -> tuple[dict, bytes]:
    paragraphs = [p.strip() for p in re.split(r"\n{1,2}", text) if p.strip()]
    if not paragraphs:
        raise ValueError("Could not extract readable text from document.")

    top_keywords = _extract_keywords(text)
    semantic_points = semantic_key_points(text, max_points=10)
    chunks = paragraphs[:24]
    grouped = [chunks[i : i + 4] for i in range(0, len(chunks), 4)]

    slides = [{"title": "Executive Snapshot", "bullets": top_keywords or ["No keywords extracted"]}]
    slides.append({"title": "AI Key Points", "bullets": semantic_points[:8] or ["No key points extracted"]})
    for idx, group in enumerate(grouped[:8], start=1):
        slides.append({"title": f"Section {idx}", "bullets": group})

    presentation = Presentation()
    title_layout = presentation.slide_layouts[0]
    content_layout = presentation.slide_layouts[1]

    cover = presentation.slides.add_slide(title_layout)
    cover.shapes.title.text = f"Report Deck: {source_name}"
    cover.placeholders[1].text = "Auto-generated from uploaded document"

    for slide_data in slides:
        slide = presentation.slides.add_slide(content_layout)
        slide.shapes.title.text = slide_data["title"][:90]
        text_frame = slide.placeholders[1].text_frame
        text_frame.clear()
        for idx, bullet in enumerate(slide_data["bullets"]):
            paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
            paragraph.text = str(bullet)[:300]
            paragraph.level = 0

    output = BytesIO()
    presentation.save(output)
    output.seek(0)
    summary = {
        "source_name": source_name,
        "slides_generated": len(slides) + 1,
        "paragraphs_analyzed": len(paragraphs),
        "top_keywords": top_keywords,
        "semantic_points": semantic_points[:8],
        "generated_at": datetime.utcnow().isoformat(),
    }
    return summary, output.getvalue()
